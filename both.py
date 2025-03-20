import streamlit as st
import json
import os
from pptx import Presentation
from pptx.util import Inches
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google_auth_oauthlib.flow import InstalledAppFlow
from twilio.rest import Client
from dotenv import load_dotenv

# 🔹 Load secrets from Streamlit OR .env file for local development
if "twilio" not in st.secrets:
    load_dotenv()
    TWILIO_ACCOUNT_SID = os.getenv("TWILIO_ACCOUNT_SID")
    TWILIO_AUTH_TOKEN = os.getenv("TWILIO_AUTH_TOKEN")
    TWILIO_PHONE_NUMBER = os.getenv("TWILIO_PHONE_NUMBER")
    RECIPIENT_PHONE_NUMBER = os.getenv("RECIPIENT_PHONE_NUMBER")

    GOOGLE_DRIVE_CREDENTIALS = json.loads(os.getenv("GOOGLE_DRIVE_CREDENTIALS"))
else:
    TWILIO_ACCOUNT_SID = st.secrets["twilio"]["TWILIO_ACCOUNT_SID"]
    TWILIO_AUTH_TOKEN = st.secrets["twilio"]["TWILIO_AUTH_TOKEN"]
    TWILIO_PHONE_NUMBER = st.secrets["twilio"]["TWILIO_PHONE_NUMBER"]
    RECIPIENT_PHONE_NUMBER = st.secrets["twilio"]["RECIPIENT_PHONE_NUMBER"]

    GOOGLE_DRIVE_CREDENTIALS = json.loads(st.secrets["google_drive"]["credentials"])

# 🔹 Authenticate Google Drive
SCOPES = ['https://www.googleapis.com/auth/drive.file']

def authenticate_drive():
    creds = None
    flow = InstalledAppFlow.from_client_config(GOOGLE_DRIVE_CREDENTIALS, SCOPES)
    creds = flow.run_console()

    return creds

# 🔹 Upload file to Google Drive
def upload_to_drive(file_path, file_name, creds):
    service = build('drive', 'v3', credentials=creds)
    media = MediaFileUpload(file_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    file_metadata = {'name': file_name}
    file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()

    file_id = file.get('id')
    service.permissions().create(
        fileId=file_id,
        body={"role": "reader", "type": "anyone"}
    ).execute()

    return f"https://drive.google.com/file/d/{file_id}/view?usp=sharing"

# 🔹 Send WhatsApp message using Twilio
def send_whatsapp_message(message):
    client = Client(TWILIO_ACCOUNT_SID, TWILIO_AUTH_TOKEN)
    client.messages.create(body=message, from_=TWILIO_PHONE_NUMBER, to=RECIPIENT_PHONE_NUMBER)

# 🔹 Append slide to PPT
def append_to_ppt(title, description, img1, img2, save_path):
    prs = Presentation(save_path) if os.path.exists(save_path) else Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    slide.shapes.title.text = title if slide.shapes.title else ""

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    textbox.text = description

    img1_path = f"temp1_{img1.name}"
    img2_path = f"temp2_{img2.name}"

    with open(img1_path, "wb") as f: f.write(img1.read())
    with open(img2_path, "wb") as f: f.write(img2.read())

    slide.shapes.add_picture(img1_path, Inches(0.5), Inches(3), width=Inches(4))
    slide.shapes.add_picture(img2_path, Inches(5), Inches(3), width=Inches(4))

    prs.save(save_path)
    os.remove(img1_path)
    os.remove(img2_path)
    return save_path

# 🔹 Streamlit UI
st.set_page_config(page_title=" Mobile PPT Generator", layout="centered")
st.title("📱 Automated PPT Generator")

title = st.text_input("Enter Slide Title:")
description = st.text_area("Enter Slide Description:")

img1 = st.file_uploader(" Capture or Upload First Image", ["png", "jpg", "jpeg"])
img2 = st.file_uploader(" Capture or Upload Second Image", ["png", "jpg", "jpeg"])

if st.button("✨ Generate PPT & Send via WhatsApp"):
    if not (title and description and img1 and img2):
        st.error(" All inputs and images are required!")
    else:
        creds = authenticate_drive()
        append_to_ppt(title, description, img1, img2, "generated_presentation.pptx")
        drive_link = upload_to_drive("generated_presentation.pptx", "generated_presentation.pptx", creds)

        send_whatsapp_message(f"PPT Ready: {drive_link}")

        st.success(" PPT created, uploaded & WhatsApp sent!")
        st.markdown(f"[View PPT]({drive_link})")
