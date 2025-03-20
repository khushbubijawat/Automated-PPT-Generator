# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches
# import os
# from googleapiclient.discovery import build
# from googleapiclient.http import MediaFileUpload
# from google_auth_oauthlib.flow import InstalledAppFlow
# import pickle
# from twilio.rest import Client

# # Configuration (update with your credentials)
# TWILIO_ACCOUNT_SID = 'ACa5421107e69fe1f668e0c624178753b4'
# TWILIO_AUTH_TOKEN = '937c440f7a4ed3591b122b4942c1b601'
# TWILIO_PHONE_NUMBER = 'whatsapp:+14155238886'
# RECIPIENT_PHONE_NUMBER = 'whatsapp:+919509920935'

# SCOPES = ['https://www.googleapis.com/auth/drive.file']
# CREDENTIALS_FILE = 'C:\\Users\\ii207\\Desktop\\Image_text_audio to ppt\\whatsup_ppt_agent\\client_secret_47430215132-u9m5h511rp2reos3ntsj73kcn8ifkja6.apps.googleusercontent.com.json'
# PPT_FILE_PATH = "mobile_generated_presentation.pptx"

# # Authenticate with Google Drive
# def authenticate_drive():
#     creds = None
#     if os.path.exists('token.pickle'):
#         with open('token.pickle', 'rb') as token:
#             creds = pickle.load(token)
#     else:
#         flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, SCOPES)
#         creds = flow.run_local_server(port=0)
#         with open('token.pickle', 'wb') as token:
#             pickle.dump(creds, token)
#     return creds

# # Upload file to Google Drive
# def upload_to_drive(file_path, file_name, creds):
#     service = build('drive', 'v3', credentials=creds)
#     media = MediaFileUpload(file_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
#     file_metadata = {'name': file_name}
#     file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
#     file_id = file.get('id')
#     service.permissions().create(
#         fileId=file_id,
#         body={"role": "reader", "type": "anyone"}
#     ).execute()
#     return f"https://drive.google.com/file/d/{file_id}/view?usp=sharing"

# # Send WhatsApp message using Twilio
# def send_whatsapp_message(message):
#     client = Client(TWILIO_ACCOUNT_SID, TWILIO_AUTH_TOKEN)
#     client.messages.create(body=message, from_=TWILIO_PHONE_NUMBER, to=RECIPIENT_PHONE_NUMBER)

# # Append slide to PPT
# def append_to_ppt(title, description, img1, img2, save_path):
#     prs = Presentation(save_path) if os.path.exists(save_path) else Presentation()
#     slide = prs.slides.add_slide(prs.slide_layouts[5])

#     slide.shapes.title.text = title if slide.shapes.title else ""

#     textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
#     textbox.text = description

#     img1_path = f"temp1_{img1.name}"
#     img2_path = f"temp2_{img2.name}"

#     with open(img1_path, "wb") as f: f.write(img1.read())
#     with open(img2_path, "wb") as f: f.write(img2.read())

#     slide.shapes.add_picture(img1_path, Inches(0.5), Inches(3), width=Inches(4))
#     slide.shapes.add_picture(img2_path, Inches(5), Inches(3), width=Inches(4))

#     prs.save(save_path)
#     os.remove(img1_path)
#     os.remove(img2_path)
#     return save_path

# # Mobile Optimized UI
# st.set_page_config(page_title="Mobile PPT Automation", layout="centered")
# st.title("📱 Automated PPT Generator")

# title = st.text_input("Enter Slide Title:")
# description = st.text_area("Enter Slide Description:")

# img1 = st.file_uploader("📸 Capture or Upload First Image", ["png", "jpg", "jpeg"])
# img2 = st.file_uploader("📸 Capture or Upload Second Image", ["png", "jpg", "jpeg"])

# if st.button(" Generate PPT"):
#     if not (title and description and img1 and img2):
#         st.error("Please provide all inputs!")
#     else:
#         append_to_ppt(title, description, img1, img2, PPT_FILE_PATH)
#         creds = authenticate_drive()
#         link = upload_to_drive(PPT_FILE_PATH, PPT_FILE_PATH, creds)
#         send_whatsapp_message(f" PPT Ready: {link}")
#         st.success(" PPT created & WhatsApp sent!")
#         st.markdown(f"[View PPT]({link})")













import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import os
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google_auth_oauthlib.flow import InstalledAppFlow
import pickle
from twilio.rest import Client
from dotenv import load_dotenv

# Load environment variables
from dotenv import load_dotenv
load_dotenv()

# Twilio Credentials (loaded from environment variables)
TWILIO_ACCOUNT_SID = os.getenv('TWILIO_ACCOUNT_SID')
TWILIO_AUTH_TOKEN = os.getenv('TWILIO_AUTH_TOKEN')
TWILIO_PHONE_NUMBER = os.getenv('TWILIO_PHONE_NUMBER')
RECIPIENT_PHONE_NUMBER = os.getenv('RECIPIENT_PHONE_NUMBER')

# Google Drive Configuration
SCOPES = ['https://www.googleapis.com/auth/drive.file']
CREDENTIALS_FILE = 'credentials.json'
PPT_FILE_PATH = "generated_presentation.pptx"

# Authenticate with Google Drive
def authenticate_drive():
    import pickle
    creds = None
    if os.path.exists('token.pickle'):
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)
    else:
        from google_auth_oauthlib.flow import InstalledAppFlow
        flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, SCOPES)
        creds = flow.run_local_server(port=0)
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)
    return creds

# Upload to Google Drive
def upload_to_drive(file_path, file_name, creds):
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload
    
    service = build('drive', 'v3', credentials=creds)
    media = MediaFileUpload(file_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    file_metadata = {'name': file_name}
    file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()

    file_id = file.get('id')
    service.permissions().create(
        fileId=file_id,
        body={"role": "reader", "type": "anyone"}
    ).execute()

    return f"https://drive.google.com/file/d/{file_id}/view?usp=sharing"

# Send WhatsApp message using Twilio
def send_whatsapp_message(message):
    client = Client(TWILIO_ACCOUNT_SID, TWILIO_AUTH_TOKEN)
    msg = client.messages.create(
        body=message,
        from_=TWILIO_PHONE_NUMBER,
        to=RECIPIENT_PHONE_NUMBER
    )
    return msg.sid

# Append slide to PPT
def append_to_ppt(title, description, img1, img2, save_path):
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation(save_path) if os.path.exists(save_path) else Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title and Description
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1)).text_frame
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
    else:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1)).text = title

    slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).text = description

    # Save Images Temporarily
    img1_path = f"temp_{img1.name}"
    img2_path = f"temp_{img2.name}"

    with open(img1_path, "wb") as f:
        f.write(img1.getbuffer())
    with open(img2_path, "wb") as f:
        f.write(img2.read())

    # Add images
    slide.shapes.add_picture(img1_path, Inches(0.5), Inches(3), width=Inches(4))
    slide.shapes.add_picture(img2_path, Inches(5), Inches(3), width=Inches(4))

    prs.save(save_path)

    # Clean-up
    os.remove(img1_path)
    os.remove(img2_path)

    return save_path

# Streamlit UI optimized for mobile
st.set_page_config(page_title="📱 Mobile PPT Generator")
st.title("📱 PPT Automation (Mobile)")

option = st.radio("Select:", ["Create New PPT", "Append to Existing PPT"])
title = st.text_input("Slide Title:")
description = st.text_area("Slide Description:")

img1 = st.file_uploader("📷 Capture or select first image", type=["png", "jpg", "jpeg"])
img2 = st.file_uploader("📷 Capture or select second image", type=["png", "jpg", "jpeg"])

if st.button("✨ Generate PPT & Send via WhatsApp"):
    if not (title and description and img1 and img2):
        st.error("⚠️ All inputs and images are required!")
    else:
        # New PPT if required
        if option == "Create New PPT":
            Presentation().save(PPT_FILE_PATH)

        # Generate and Upload
        append_to_ppt(title, description, img1, img2, PPT_FILE_PATH)
        creds = authenticate_drive()
        drive_link = upload_to_drive(PPT_FILE_PATH, PPT_FILE_PATH, creds)

        # Send WhatsApp (fully automatic via Twilio)
        send_whatsapp_message(f"✅ PPT is ready! Download here: {drive_link}")

        st.success("✅ PPT created, uploaded, WhatsApp sent!")
        st.markdown(f"[View PPT]({drive_link})")
